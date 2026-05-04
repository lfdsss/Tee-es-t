"""SNB Mission Hunter — PDF and PPTX document generation.
NE JAMAIS mentionner S&B Consulting dans les documents generes.
"""

import io
from datetime import datetime
from typing import Dict, Any, List

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE_ACCENT = (37, 99, 235)
DARK_BLUE = (30, 58, 95)
WHITE = (255, 255, 255)
LIGHT_GRAY = (245, 247, 250)
DARK_TEXT = (17, 17, 17)
GRAY_TEXT = (100, 100, 100)

CONTACT_EMAIL = "bp.thevenot@gmail.com"
CONTACT_PHONE = "06 86 50 43 79"
SIRET = "849 022 058"
CONSULTANT_NAME = "Baptiste Thevenot"
CONSULTANT_TITLE = "Consultant Web & IA"
FOOTER_TEXT = f"Baptiste Thevenot — SIRET {SIRET} — TVA non applicable art. 293B du CGI"


def _get_deliverable_text(d) -> str:
    if isinstance(d, dict):
        name = d.get("name", "")
        desc = d.get("description", "")
        fmt = d.get("format", "")
        parts = [name]
        if fmt:
            parts.append(f"({fmt})")
        if desc:
            parts.append(f": {desc}")
        return " ".join(parts)
    return str(d)


def _get_phases(pkg: Dict[str, Any]) -> List[Dict[str, Any]]:
    return pkg.get("phases", [])


class ProposalPDF(FPDF):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.set_auto_page_break(auto=True, margin=25)

    def header(self):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(*DARK_BLUE)
        self.cell(0, 6, f"{CONSULTANT_NAME} — {CONSULTANT_TITLE}", align="L")
        self.ln(4)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(*GRAY_TEXT)
        self.cell(0, 4, f"{CONTACT_EMAIL} | {CONTACT_PHONE} | SIRET {SIRET}", align="L")
        self.ln(6)
        self.set_draw_color(*BLUE_ACCENT)
        self.set_line_width(0.5)
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(6)

    def footer(self):
        self.set_y(-20)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(*GRAY_TEXT)
        self.cell(0, 5, FOOTER_TEXT, align="C")
        self.ln(3)
        self.cell(0, 5, f"Page {self.page_no()}/{{nb}}", align="C")

    def section_title(self, title: str):
        self.ln(4)
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(*DARK_BLUE)
        self.cell(0, 8, title)
        self.ln(3)
        self.set_draw_color(*BLUE_ACCENT)
        self.set_line_width(0.3)
        self.line(10, self.get_y(), 80, self.get_y())
        self.ln(5)

    def body_text(self, text: str):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*DARK_TEXT)
        self.multi_cell(0, 5, text)
        self.ln(2)

    def bullet(self, text: str):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*DARK_TEXT)
        self.cell(6, 5, chr(8226))
        self.multi_cell(0, 5, text)
        self.ln(1)


class PDFGenerator:
    def generate(self, package_json: Dict[str, Any], mission_data: Dict[str, Any]) -> bytes:
        pdf = ProposalPDF()
        pdf.alias_nb_pages()
        pdf.add_page()

        today_str = datetime.now().strftime("%d/%m/%Y")
        title = mission_data.get("title", "Mission")
        client = mission_data.get("company", "Client")

        pdf.set_font("Helvetica", "B", 20)
        pdf.set_text_color(*DARK_BLUE)
        pdf.cell(0, 12, "PROPOSITION DE MISSION", align="C")
        pdf.ln(10)

        pdf.set_font("Helvetica", "", 12)
        pdf.set_text_color(*DARK_TEXT)
        pdf.cell(0, 7, title, align="C")
        pdf.ln(6)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*GRAY_TEXT)
        pdf.cell(0, 6, f"Client : {client}", align="C")
        pdf.ln(4)
        pdf.cell(0, 6, f"Date : {today_str}", align="C")
        pdf.ln(10)

        if package_json.get("executive_summary"):
            pdf.section_title("Resume Executif")
            pdf.body_text(package_json["executive_summary"])

        if package_json.get("comprehension"):
            pdf.section_title("Comprehension du Besoin")
            pdf.body_text(package_json["comprehension"])

        if package_json.get("approach"):
            pdf.section_title("Approche")
            pdf.body_text(package_json["approach"])

        if package_json.get("methodology"):
            pdf.body_text(package_json["methodology"])

        phases = _get_phases(package_json)
        if phases:
            pdf.section_title("Phases du Projet")
            for i, phase in enumerate(phases, 1):
                pdf.set_font("Helvetica", "B", 11)
                pdf.set_text_color(*BLUE_ACCENT)
                name = phase.get("name", f"Phase {i}")
                duration = phase.get("duration", "")
                pdf.cell(0, 7, f"{name} — {duration}")
                pdf.ln(5)

                if phase.get("objective"):
                    pdf.set_font("Helvetica", "I", 10)
                    pdf.set_text_color(*GRAY_TEXT)
                    pdf.cell(0, 5, f"Objectif : {phase['objective']}")
                    pdf.ln(5)

                for task in phase.get("tasks", []):
                    pdf.bullet(task)

                if phase.get("deliverable"):
                    pdf.set_font("Helvetica", "I", 9)
                    pdf.set_text_color(*GRAY_TEXT)
                    pdf.cell(0, 5, f"Livrable : {phase['deliverable']}")
                    pdf.ln(4)

                if phase.get("tools"):
                    pdf.set_font("Helvetica", "I", 9)
                    pdf.set_text_color(*GRAY_TEXT)
                    pdf.cell(0, 5, f"Outils : {', '.join(phase['tools'])}")
                    pdf.ln(4)

                pdf.ln(3)

        if package_json.get("timeline"):
            pdf.section_title("Planning")
            pdf.body_text(f"Duree totale estimee : {package_json['timeline']}")

        deliverables = package_json.get("deliverables", [])
        if deliverables:
            pdf.section_title("Livrables")
            for d in deliverables:
                pdf.bullet(_get_deliverable_text(d))

        pricing = package_json.get("pricing", {})
        if pricing:
            pdf.section_title("Tarification")

            detail = pricing.get("detail", [])
            if detail:
                pdf.set_font("Helvetica", "B", 9)
                pdf.set_fill_color(*LIGHT_GRAY)
                pdf.set_text_color(*DARK_BLUE)
                col_w = [80, 30, 50]
                pdf.cell(col_w[0], 7, "Poste", border=1, fill=True)
                pdf.cell(col_w[1], 7, "Jours", border=1, fill=True, align="C")
                pdf.cell(col_w[2], 7, "Montant", border=1, fill=True, align="R")
                pdf.ln()

                pdf.set_font("Helvetica", "", 9)
                pdf.set_text_color(*DARK_TEXT)
                for item in detail:
                    pdf.cell(col_w[0], 6, str(item.get("item", "")), border=1)
                    pdf.cell(col_w[1], 6, str(item.get("days", "")), border=1, align="C")
                    pdf.cell(col_w[2], 6, str(item.get("amount", "")), border=1, align="R")
                    pdf.ln()

                pdf.set_font("Helvetica", "B", 10)
                pdf.set_text_color(*DARK_BLUE)
                pdf.cell(col_w[0] + col_w[1], 7, "TOTAL", border=1, fill=True, align="R")
                pdf.cell(col_w[2], 7, str(pricing.get("amount", "")), border=1, fill=True, align="R")
                pdf.ln(8)

            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(*DARK_TEXT)
            pdf.body_text(f"Modele : {pricing.get('model', '')}")
            pdf.body_text(f"Montant total : {pricing.get('amount', '')}")
            pdf.body_text(f"Conditions de paiement : {pricing.get('payment', '')}")

        guarantees = package_json.get("guarantees", [])
        if guarantees:
            pdf.section_title("Garanties")
            for g in guarantees:
                pdf.bullet(g)

        if package_json.get("next_step"):
            pdf.section_title("Prochaine Etape")
            pdf.body_text(package_json["next_step"])

        pdf.ln(10)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(*GRAY_TEXT)
        pdf.cell(0, 5, "Bon pour accord", align="L")
        pdf.ln(15)
        pdf.set_draw_color(*GRAY_TEXT)
        pdf.line(10, pdf.get_y(), 70, pdf.get_y())
        pdf.ln(3)
        pdf.cell(0, 5, "Signature client", align="L")

        pdf.ln(10)
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*DARK_BLUE)
        pdf.cell(0, 5, package_json.get("signature", f"{CONSULTANT_NAME}, {CONSULTANT_TITLE}"), align="R")

        buf = io.BytesIO()
        pdf.output(buf)
        return buf.getvalue()


class PPTXGenerator:
    def generate(self, package_json: Dict[str, Any], mission_data: Dict[str, Any]) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title = mission_data.get("title", "Mission")
        client = mission_data.get("company", "Client")
        today_str = datetime.now().strftime("%d/%m/%Y")
        phases = _get_phases(package_json)

        self._slide_cover(prs, title, client, today_str)
        self._slide_executive_summary(prs, package_json)
        self._slide_besoin(prs, package_json)
        self._slide_approche(prs, package_json)

        for i, phase in enumerate(phases):
            self._slide_phase(prs, phase, i + 1)

        self._slide_planning_livrables(prs, package_json)
        self._slide_investissement(prs, package_json)
        self._slide_next_steps(prs, package_json)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _add_title_bar(self, slide, title_text: str):
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.333)
        height = Inches(1.1)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)

    def _add_content_text(self, slide, text: str, left: float, top: float,
                          width: float, height: float, font_size: int = 14,
                          bold: bool = False, color=DARK_TEXT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        return tf

    def _add_bullet_list(self, slide, items: list, left: float, top: float,
                         width: float, height: float, font_size: int = 13):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*DARK_TEXT)
            p.space_after = Pt(6)
        return tf

    def _slide_cover(self, prs, title: str, client: str, date: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
        bg_shape.line.fill.background()

        accent = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*BLUE_ACCENT)
        accent.line.fill.background()

        self._add_content_text(slide, "PROPOSITION DE MISSION", 1.5, 1.5, 10, 0.6, 18, False, WHITE)
        self._add_content_text(slide, title, 1.5, 2.2, 10, 1.0, 32, True, WHITE)
        self._add_content_text(slide, f"Client : {client}", 1.5, 3.8, 10, 0.5, 16, False, (180, 200, 230))
        self._add_content_text(slide, date, 1.5, 4.4, 10, 0.4, 14, False, (150, 170, 200))

        self._add_content_text(slide, f"{CONSULTANT_NAME} — {CONSULTANT_TITLE}", 1.5, 6.0, 10, 0.4, 14, False, (180, 200, 230))
        self._add_content_text(slide, f"{CONTACT_EMAIL} | {CONTACT_PHONE}", 1.5, 6.4, 10, 0.3, 11, False, (140, 160, 190))

    def _slide_executive_summary(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Resume Executif")

        summary = pkg.get("executive_summary", pkg.get("intro", ""))
        self._add_content_text(slide, summary, 1.0, 1.5, 11, 1.5, 16, False, DARK_TEXT)

        kpis = pkg.get("kpis", [])
        if kpis:
            self._add_content_text(slide, "Indicateurs cles :", 1.0, 3.5, 11, 0.4, 14, True, DARK_BLUE)
            self._add_bullet_list(slide, kpis, 1.0, 4.0, 11, 2.0, 13)

        tools = pkg.get("tools_used", [])
        if tools:
            self._add_content_text(slide, "Technologies :", 1.0, 5.5, 11, 0.4, 14, True, DARK_BLUE)
            self._add_content_text(slide, ", ".join(tools), 1.0, 5.9, 11, 0.5, 13, False, GRAY_TEXT)

    def _slide_besoin(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Votre Besoin")

        comprehension = pkg.get("comprehension", "")
        self._add_content_text(slide, comprehension, 1.0, 1.5, 11, 2.5, 15, False, DARK_TEXT)

        self._add_content_text(slide, "Ce que nous avons compris :", 1.0, 4.5, 11, 0.4, 14, True, DARK_BLUE)
        intro = pkg.get("intro", "")
        self._add_content_text(slide, intro, 1.0, 5.0, 11, 1.5, 13, False, GRAY_TEXT)

    def _slide_approche(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Notre Approche")

        approach = pkg.get("approach", "")
        self._add_content_text(slide, approach, 1.0, 1.5, 11, 1.5, 15, False, DARK_TEXT)

        methodology = pkg.get("methodology", "")
        if methodology:
            self._add_content_text(slide, "Methodologie", 1.0, 3.5, 11, 0.4, 14, True, DARK_BLUE)
            self._add_content_text(slide, methodology, 1.0, 4.0, 11, 1.5, 13, False, DARK_TEXT)

        tools = pkg.get("tools_used", [])
        if tools:
            self._add_content_text(slide, "Outils & Technologies", 1.0, 5.8, 11, 0.4, 14, True, DARK_BLUE)
            self._add_bullet_list(slide, tools, 1.0, 6.2, 11, 1.0, 12)

    def _slide_phase(self, prs, phase: Dict[str, Any], index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        name = phase.get("name", f"Phase {index}")
        self._add_title_bar(slide, name)

        if phase.get("objective"):
            self._add_content_text(slide, f"Objectif : {phase['objective']}", 1.0, 1.4, 11, 0.5, 14, True, BLUE_ACCENT)

        y_pos = 2.2 if phase.get("objective") else 1.6
        duration = phase.get("duration", "")
        self._add_content_text(slide, f"Duree : {duration}", 1.0, y_pos, 5, 0.4, 13, False, GRAY_TEXT)

        tasks = phase.get("tasks", [])
        if tasks:
            self._add_content_text(slide, "Taches :", 1.0, y_pos + 0.6, 11, 0.3, 13, True, DARK_BLUE)
            self._add_bullet_list(slide, tasks, 1.0, y_pos + 1.0, 11, 2.5, 12)

        deliverable = phase.get("deliverable", "")
        if deliverable:
            self._add_content_text(slide, "Livrable :", 1.0, 5.5, 5, 0.3, 13, True, DARK_BLUE)
            self._add_content_text(slide, deliverable, 1.0, 5.9, 11, 0.5, 12, False, DARK_TEXT)

        tools = phase.get("tools", [])
        if tools:
            self._add_content_text(slide, "Outils :", 1.0, 6.5, 5, 0.3, 13, True, DARK_BLUE)
            self._add_content_text(slide, ", ".join(tools), 2.5, 6.5, 9, 0.3, 12, False, GRAY_TEXT)

    def _slide_planning_livrables(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Planning & Livrables")

        timeline = pkg.get("timeline", "")
        self._add_content_text(slide, f"Duree totale : {timeline}", 1.0, 1.5, 11, 0.4, 15, True, DARK_BLUE)

        phases = _get_phases(pkg)
        if phases:
            self._add_content_text(slide, "Phases :", 1.0, 2.2, 11, 0.3, 13, True, DARK_BLUE)
            phase_texts = [f"{p.get('name', f'Phase {i+1}')} — {p.get('duration', '')}" for i, p in enumerate(phases)]
            self._add_bullet_list(slide, phase_texts, 1.0, 2.6, 11, 1.5, 12)

        deliverables = pkg.get("deliverables", [])
        if deliverables:
            y_start = 4.5
            self._add_content_text(slide, "Livrables finaux :", 1.0, y_start, 11, 0.3, 13, True, DARK_BLUE)
            deliv_texts = [_get_deliverable_text(d) for d in deliverables]
            self._add_bullet_list(slide, deliv_texts, 1.0, y_start + 0.4, 11, 2.5, 12)

    def _slide_investissement(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Investissement")

        pricing = pkg.get("pricing", {})

        detail = pricing.get("detail", [])
        if detail:
            y = 1.6
            headers = ["Poste", "Jours", "Montant"]
            col_x = [1.0, 7.0, 9.0]
            col_w = [6.0, 2.0, 3.0]

            for ci, header in enumerate(headers):
                shape = slide.shapes.add_shape(1, Inches(col_x[ci]), Inches(y), Inches(col_w[ci]), Inches(0.4))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
                shape.line.fill.background()
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = header
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*WHITE)
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

            for ri, item in enumerate(detail):
                row_y = y + 0.45 + ri * 0.4
                bg_color = LIGHT_GRAY if ri % 2 == 0 else WHITE
                values = [str(item.get("item", "")), str(item.get("days", "")), str(item.get("amount", ""))]
                for ci, val in enumerate(values):
                    shape = slide.shapes.add_shape(1, Inches(col_x[ci]), Inches(row_y), Inches(col_w[ci]), Inches(0.38))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*bg_color)
                    shape.line.color.rgb = RGBColor(220, 220, 220)
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = val
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(*DARK_TEXT)
                    p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

            total_y = y + 0.45 + len(detail) * 0.4
            shape = slide.shapes.add_shape(1, Inches(col_x[0]), Inches(total_y), Inches(col_w[0] + col_w[1]), Inches(0.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*BLUE_ACCENT)
            shape.line.fill.background()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "TOTAL"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*WHITE)
            p.alignment = PP_ALIGN.RIGHT

            shape = slide.shapes.add_shape(1, Inches(col_x[2]), Inches(total_y), Inches(col_w[2]), Inches(0.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*BLUE_ACCENT)
            shape.line.fill.background()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = str(pricing.get("amount", ""))
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*WHITE)
            p.alignment = PP_ALIGN.CENTER

            info_y = total_y + 0.8
        else:
            self._add_content_text(slide, f"Montant : {pricing.get('amount', '')}", 1.0, 2.0, 11, 0.5, 18, True, DARK_BLUE)
            info_y = 3.0

        model = pricing.get("model", "")
        payment = pricing.get("payment", "")
        self._add_content_text(slide, f"Modele : {model}", 1.0, info_y, 11, 0.4, 13, False, DARK_TEXT)
        self._add_content_text(slide, f"Paiement : {payment}", 1.0, info_y + 0.5, 11, 0.4, 13, False, DARK_TEXT)
        self._add_content_text(slide, "TVA non applicable — art. 293B du CGI", 1.0, info_y + 1.2, 11, 0.3, 10, False, GRAY_TEXT)

    def _slide_next_steps(self, prs, pkg: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_bar(slide, "Prochaines Etapes")

        next_step = pkg.get("next_step", "")
        self._add_content_text(slide, next_step, 1.0, 1.5, 11, 1.5, 16, False, DARK_TEXT)

        guarantees = pkg.get("guarantees", [])
        if guarantees:
            self._add_content_text(slide, "Nos Garanties", 1.0, 3.5, 11, 0.4, 14, True, DARK_BLUE)
            self._add_bullet_list(slide, guarantees, 1.0, 4.0, 11, 1.5, 13)

        self._add_content_text(slide, "Contact", 1.0, 5.5, 11, 0.4, 14, True, DARK_BLUE)
        self._add_content_text(slide, f"{CONSULTANT_NAME} — {CONSULTANT_TITLE}", 1.0, 5.9, 11, 0.4, 13, True, DARK_TEXT)
        self._add_content_text(slide, f"{CONTACT_EMAIL} | {CONTACT_PHONE}", 1.0, 6.3, 11, 0.3, 12, False, GRAY_TEXT)
        self._add_content_text(slide, f"SIRET {SIRET}", 1.0, 6.6, 11, 0.3, 10, False, GRAY_TEXT)

        signature = pkg.get("signature", f"{CONSULTANT_NAME}, {CONSULTANT_TITLE}")
        self._add_content_text(slide, signature, 7.0, 6.6, 5, 0.3, 11, True, DARK_BLUE)
